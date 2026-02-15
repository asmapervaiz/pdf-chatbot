"""
Build PDF Chatbot presentation (PowerPoint .pptx).
Run from project root:  python scripts/build_presentation.py
Requires: pip install python-pptx
"""

import os
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("Install python-pptx: pip install python-pptx")
    sys.exit(1)

# Project root for output path
ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
os.chdir(ROOT)


def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide


def add_content_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(bullets):
        p = tf.add_paragraph()
        p.text = line
        p.space_after = Pt(6)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(
        prs,
        "PDF Chatbot",
        "AI/ML Document Q&A — AppLab Assignment\nPython • FastAPI • Streamlit • RAG",
    )

    # Slide 2: Objective
    add_content_slide(prs, "Objective", [
        "Build a Python app that answers user questions from uploaded PDFs.",
        "Document Upload API: upload PDFs, extract and index text.",
        "Chatbot API: ask questions, get answers using NLP/AI (RAG).",
        "User interface: upload + chat (Streamlit).",
        "Deploy with Docker; maintain clear, documented code.",
    ])

    # Slide 3: Architecture overview
    add_content_slide(prs, "Application Design & Architecture", [
        "Backend: FastAPI (app.main) — REST API only.",
        "Document flow: Upload PDF → PyMuPDF extract text → chunk by paragraphs → embed (OpenAI or sentence-transformers) → store in ChromaDB.",
        "Chat flow: Question → embed → similarity search (top-k chunks) → build context → LLM (OpenAI or FLAN-T5) → answer + sources.",
        "UI: Streamlit (streamlit_app.py) — calls backend APIs (upload, clear, chat/ask).",
        "Single command: python run.py starts backend (thread) + Streamlit (subprocess).",
    ])

    # Slide 4: NLP & AI model selection
    add_content_slide(prs, "NLP & AI Model Selection", [
        "Embeddings: OpenAI text-embedding-3-small (when API key set) or sentence-transformers/all-MiniLM-L6-v2 (local, 384-dim). Separate ChromaDB collection per dim to avoid mismatch.",
        "Vector store: ChromaDB — persistent, simple API, good for RAG.",
        "Retrieval: top-k similarity + optional query expansion (key phrases from question) for policy-style Q&A.",
        "Answer generation: OpenAI gpt-3.5-turbo (if key set) or HuggingFace google/flan-t5-small (local). Prompt instructs model to answer from context only.",
    ])

    # Slide 5: API overview
    add_content_slide(prs, "API Overview", [
        "POST /documents/upload — multipart PDF; returns filename, pages_processed, chunks_indexed.",
        "GET /documents/status — returns chunks_indexed (for UI/debug).",
        "POST /documents/clear — clear vector store.",
        "POST /chat/ask — JSON {question}; returns {answer, sources}.",
        "GET /health — health check. GET /docs — OpenAPI docs.",
    ])

    # Slide 6: Deployment & usage
    add_content_slide(prs, "Deployment & Usage", [
        "Local: python run.py (backend + Streamlit). Open http://localhost:8501 for UI, http://localhost:8000/docs for API.",
        "Optional: set OPENAI_API_KEY in .env for OpenAI embeddings and chat.",
        "Docker: docker-compose up --build. Volumes: ./uploads, ./data for persistence.",
        "If you switch embedding provider, call POST /documents/clear then re-upload PDFs.",
    ])

    # Slide 7: Project structure
    add_content_slide(prs, "Project Structure", [
        "app/main.py — FastAPI app, routes, lifespan.",
        "app/api/ — documents.py (upload, status, clear), chat.py (ask).",
        "app/services/ — pdf_service (extract, chunk), embeddings_service (ChromaDB), chat_service (RAG + LLM).",
        "app/config.py — settings; app/models/schemas.py — Pydantic models.",
        "streamlit_app.py — Streamlit UI. run.py — single-command launcher.",
    ])

    # Slide 8: Live demo / summary
    add_content_slide(prs, "Live Demo & Summary", [
        "Demo: Run python run.py → open http://localhost:8501 → upload a PDF → ask questions (e.g. leave policy, password length).",
        "Features: paragraph-aware chunking, shared embedding service (upload + chat same index), clear index, optional OpenAI.",
        "Deliverables: GitHub repo, Dockerfile, README with setup/API docs, this presentation.",
    ])

    out_path = os.path.join(ROOT, "PDF_Chatbot_Presentation.pptx")
    prs.save(out_path)
    print("Saved:", out_path)


if __name__ == "__main__":
    main()
